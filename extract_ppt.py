from pptx import Presentation
import os

def extract_text_from_ppt(ppt_path):
    print(f"--- Extracting from {os.path.basename(ppt_path)} ---")
    try:
        prs = Presentation(ppt_path)
        text_content = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"Slide {i+1}:")
            
            # Extract title
            if slide.shapes.title:
                slide_text.append(f"Title: {slide.shapes.title.text}")
            
            # Extract all other text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text != slide.shapes.title.text:
                    slide_text.append(shape.text)
            
            text_content.append("\n".join(slide_text))
        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error reading {ppt_path}: {str(e)}"

if __name__ == "__main__":
    ppt_files = ["C14_Review.pptx", "C14_PPT2.pptx", "C14_Review[1].pptx"]
    
    with open("ppt_extracted_content.txt", "w", encoding="utf-8") as f:
        for ppt_file in ppt_files:
            if os.path.exists(ppt_file):
                content = extract_text_from_ppt(ppt_file)
                f.write(content)
                f.write("\n\n" + "="*50 + "\n\n")
            else:
                print(f"File not found: {ppt_file}")
                
    print("Extraction complete. Saved to ppt_extracted_content.txt")
